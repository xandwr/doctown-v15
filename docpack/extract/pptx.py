"""PPTX document extraction using python-pptx."""

from __future__ import annotations

import io

from .base import ExtractedDocument, ExtractedImage


def extract_pptx(data: bytes) -> ExtractedDocument:
    """
    Extract text and images from PPTX bytes.

    Args:
        data: Raw PPTX file bytes

    Returns:
        ExtractedDocument with text content and embedded images
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(data))
    text_parts: list[str] = []
    images: list[ExtractedImage] = []
    global_image_index = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text_parts: list[str] = []
        current_title: str | None = None

        # Get slide title if available
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                current_title = title_text
                slide_text_parts.append(f"# {title_text}")

        # Extract text from all shapes
        for shape in slide.shapes:
            # Handle text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text and text != current_title:
                        slide_text_parts.append(text)

            # Handle tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = []
                    for cell in row.cells:
                        cells.append(cell.text.strip())
                    if any(cells):
                        slide_text_parts.append(" | ".join(cells))

            # Handle images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    img_data = image.blob
                    content_type = image.content_type

                    # Determine format
                    if "png" in content_type:
                        fmt = "png"
                    elif "jpeg" in content_type or "jpg" in content_type:
                        fmt = "jpeg"
                    elif "gif" in content_type:
                        fmt = "gif"
                    else:
                        fmt = content_type.split("/")[-1] if "/" in content_type else "unknown"

                    # Get dimensions
                    width, height = 0, 0
                    try:
                        from PIL import Image

                        pil_img = Image.open(io.BytesIO(img_data))
                        width, height = pil_img.size
                    except Exception:
                        pass

                    images.append(
                        ExtractedImage(
                            data=img_data,
                            format=fmt,
                            width=width,
                            height=height,
                            page_number=slide_num,
                            image_index=global_image_index,
                            context=current_title,
                        )
                    )
                    global_image_index += 1
                except Exception:
                    continue

        if slide_text_parts:
            text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text_parts))

    # Get metadata from core properties
    metadata: dict[str, str] = {}
    try:
        if prs.core_properties:
            if prs.core_properties.title:
                metadata["title"] = prs.core_properties.title
            if prs.core_properties.author:
                metadata["author"] = prs.core_properties.author
    except Exception:
        pass

    return ExtractedDocument(
        text="\n\n".join(text_parts),
        images=images,
        page_count=len(prs.slides),
        metadata=metadata,
    )
